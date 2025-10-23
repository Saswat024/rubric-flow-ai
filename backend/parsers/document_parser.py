"""
Document parser for extracting text from various file formats
Supports: .docx, .pptx, .pdf, .txt
"""
import io
import base64
from typing import Dict, Any
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader


async def parse_document(file_base64: str, file_type: str) -> Dict[str, Any]:
    """
    Parse document and extract text content
    
    Args:
        file_base64: Base64 encoded file content
        file_type: File extension (.docx, .pptx, .pdf, .txt)
    
    Returns:
        Dictionary containing extracted text and metadata
    """
    try:
        # Decode base64 content
        file_content = base64.b64decode(file_base64.split(',')[1] if ',' in file_base64 else file_base64)
        file_stream = io.BytesIO(file_content)
        
        extracted_text = ""
        metadata = {
            "type": file_type,
            "pages": 0,
            "slides": 0,
            "paragraphs": 0
        }
        
        if file_type == ".docx":
            extracted_text, metadata = await parse_docx(file_stream)
        elif file_type == ".pptx":
            extracted_text, metadata = await parse_pptx(file_stream)
        elif file_type == ".pdf":
            extracted_text, metadata = await parse_pdf(file_stream)
        elif file_type == ".txt":
            extracted_text = file_content.decode('utf-8')
            metadata["paragraphs"] = len(extracted_text.split('\n\n'))
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return {
            "text": extracted_text,
            "metadata": metadata,
            "success": True
        }
    
    except Exception as e:
        print(f"Error parsing document: {str(e)}")
        return {
            "text": "",
            "metadata": {},
            "success": False,
            "error": str(e)
        }


async def parse_docx(file_stream: io.BytesIO) -> tuple[str, Dict[str, Any]]:
    """Parse DOCX file and extract text"""
    doc = Document(file_stream)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    text = "\n\n".join(paragraphs)
    
    metadata = {
        "type": ".docx",
        "paragraphs": len(paragraphs),
        "pages": len(doc.sections)
    }
    
    return text, metadata


async def parse_pptx(file_stream: io.BytesIO) -> tuple[str, Dict[str, Any]]:
    """Parse PPTX file and extract text from slides"""
    prs = Presentation(file_stream)
    slide_texts = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        if slide_content:
            slide_texts.append(f"Slide {i + 1}:\n" + "\n".join(slide_content))
    
    text = "\n\n".join(slide_texts)
    
    metadata = {
        "type": ".pptx",
        "slides": len(prs.slides),
        "pages": len(prs.slides)
    }
    
    return text, metadata


async def parse_pdf(file_stream: io.BytesIO) -> tuple[str, Dict[str, Any]]:
    """Parse PDF file and extract text"""
    pdf_reader = PdfReader(file_stream)
    page_texts = []
    
    for i, page in enumerate(pdf_reader.pages):
        page_text = page.extract_text()
        if page_text.strip():
            page_texts.append(f"Page {i + 1}:\n{page_text}")
    
    text = "\n\n".join(page_texts)
    
    metadata = {
        "type": ".pdf",
        "pages": len(pdf_reader.pages)
    }
    
    return text, metadata
